"""Office extraction: .docx (headings preserved), .pptx (per slide), .xlsx (per sheet).

Word documents come out as Markdown so the heading-aware chunker can work on them —
that is why ``extract_docx`` returns ``KIND_MARKDOWN`` rather than plain text.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from . import KIND_MARKDOWN, KIND_PAGED, KIND_TABULAR, Block, Document

MAX_SHEET_ROWS = 20_000
ROWS_PER_BLOCK = 40


def extract_docx(path: Path, **_: Any) -> Document:
    """Word to Markdown: real headings become ``#`` levels, tables become pipe rows."""
    import docx  # python-docx, from the 'documents' extra

    source = docx.Document(str(path))
    lines: list[str] = []
    for paragraph in source.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = "".join(c for c in style if c.isdigit())
            lines.append(f"{'#' * min(int(level or 1), 6)} {text}")
        elif style.startswith("list") or style.startswith("bullet"):
            lines.append(f"- {text}")
        else:
            lines.append(text)

    for number, table in enumerate(source.tables, start=1):
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        lines.append(f"\n**Table {number}**\n")
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("|" + "|".join(" --- " for _ in rows[0]) + "|")
        lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])

    props = source.core_properties
    title = (props.title or "").strip() or path.stem
    return Document(
        path=str(path), kind=KIND_MARKDOWN, title=title, text="\n\n".join(lines),
        meta={
            "author": (props.author or "").strip(),
            "paragraphs": len(source.paragraphs),
            "tables": len(source.tables),
        },
    )


def extract_pptx(path: Path, **_: Any) -> Document:
    """One block per slide, so hits cite a slide number. Speaker notes are included."""
    from pptx import Presentation  # python-pptx, from the 'documents' extra

    deck = Presentation(str(path))
    doc = Document(path=str(path), kind=KIND_PAGED, title=path.stem)
    for number, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")
        if parts:
            doc.blocks.append(
                Block(text="\n".join(parts), kind="slide", anchor={"slide": number})
            )
    doc.meta = {"slide_count": len(deck.slides)}
    return doc


def extract_xlsx(path: Path, rows_per_block: int = ROWS_PER_BLOCK, **_: Any) -> Document:
    """Row groups per sheet, header repeated in each block so every block stands alone."""
    from openpyxl import load_workbook  # from the 'documents' extra

    book = load_workbook(str(path), read_only=True, data_only=True)
    doc = Document(path=str(path), kind=KIND_TABULAR, title=path.stem)
    sheets: list[str] = []
    try:
        for sheet in book.worksheets:
            sheets.append(sheet.title)
            header: list[str] = []
            batch: list[str] = []
            start = 2
            for number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if number > MAX_SHEET_ROWS:
                    doc.meta["truncated_sheets"] = doc.meta.get("truncated_sheets", []) + [sheet.title]
                    break
                cells = ["" if value is None else str(value).strip() for value in row]
                if not any(cells):
                    continue
                if not header:
                    header = cells
                    continue
                batch.append(" | ".join(cells))
                if len(batch) >= rows_per_block:
                    doc.blocks.append(_sheet_block(sheet.title, header, batch, start, number))
                    batch, start = [], number + 1
            if batch:
                doc.blocks.append(
                    _sheet_block(sheet.title, header, batch, start, start + len(batch) - 1)
                )
    finally:
        book.close()

    doc.meta["sheets"] = sheets
    return doc


def _sheet_block(sheet: str, header: list[str], rows: list[str], first: int, last: int) -> Block:
    head = " | ".join(header)
    body = f"Sheet '{sheet}' — columns: {head}\n" + "\n".join(rows)
    return Block(
        text=body, kind="rows",
        anchor={"sheet": sheet, "row_start": first, "row_end": last},
    )
